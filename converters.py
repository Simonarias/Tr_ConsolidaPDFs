import io
import os
import shutil
import subprocess
import tempfile
from PIL import Image
from pypdf import PdfReader, PdfWriter

# ReportLab imports
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Preformatted, PageBreak, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

# Optional office library imports
try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pptx
except ImportError:
    pptx = None

# Extensiones clasificadas por tipo
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp', '.bmp', '.tiff', '.tif', '.gif'}
TEXT_EXTENSIONS = {'.txt', '.csv', '.md', '.json', '.log', '.py', '.js', '.html', '.css', '.xml', '.yaml', '.yml'}
WORD_EXTENSIONS = {'.docx'}
EXCEL_EXTENSIONS = {'.xlsx'}
PPT_EXTENSIONS = {'.pptx'}

def _convert_with_libreoffice(input_file_path, output_dir):
    """Convierte un archivo de Office (docx, xlsx, pptx) a PDF usando LibreOffice en modo headless."""
    executables = [
        'libreoffice', 'soffice',
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe'
    ]
    cmd = None
    for exe in executables:
        if shutil.which(exe) or os.path.exists(exe):
            cmd = exe
            break
    if not cmd:
        return False
        
    try:
        res = subprocess.run(
            [cmd, '--headless', '--convert-to', 'pdf', input_file_path, '--outdir', output_dir],
            capture_output=True,
            text=True,
            timeout=120
        )
        base_name = os.path.splitext(os.path.basename(input_file_path))[0]
        generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")
        target_pdf = os.path.join(output_dir, "output.pdf")
        if os.path.exists(generated_pdf):
            if generated_pdf != target_pdf:
                if os.path.exists(target_pdf):
                    os.remove(target_pdf)
                os.rename(generated_pdf, target_pdf)
            return True
    except Exception:
        pass
    return False

def _convert_docx_with_win32com(input_docx_path, output_pdf_path):
    """Convierte Word (.docx) a PDF usando la API COM de Microsoft Word en Windows."""
    if os.name != 'nt':
        return False
    try:
        import win32com.client
        pythoncom_initialized = False
        try:
            import pythoncom
            pythoncom.CoInitialize()
            pythoncom_initialized = True
        except Exception:
            pass

        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        try:
            doc = word.Documents.Open(os.path.abspath(input_docx_path))
            doc.SaveAs(os.path.abspath(output_pdf_path), FileFormat=17) # 17 = wdFormatPDF
            doc.Close()
            return os.path.exists(output_pdf_path)
        finally:
            word.Quit()
            if pythoncom_initialized:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass
    except Exception:
        pass
    return False

def convert_image_to_pdf(file_bytes_or_stream):
    """Convierte una imagen (JPG, PNG, WEBP, etc.) a PDF en memoria."""
    img = Image.open(file_bytes_or_stream)
    if img.mode in ("RGBA", "P", "LA"):
        img = img.convert("RGB")
    pdf_buffer = io.BytesIO()
    img.save(pdf_buffer, format="PDF", resolution=100.0)
    pdf_buffer.seek(0)
    return pdf_buffer

def convert_text_to_pdf(text_str, filename="Documento"):
    """Convierte texto plano o código fuente a PDF con formato legible."""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=15,
        leading=18,
        textColor=colors.HexColor('#4F46E5'),
        spaceAfter=12
    )
    
    code_style = ParagraphStyle(
        'CodeStyle',
        parent=styles['Code'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=13,
        textColor=colors.HexColor('#1E293B')
    )
    
    story = []
    story.append(Paragraph(f"📄 {filename}", title_style))
    story.append(Spacer(1, 10))
    
    lines = text_str.splitlines()
    if not lines:
        lines = ["(Archivo de texto vacío)"]
        
    chunk_size = 60
    for i in range(0, len(lines), chunk_size):
        chunk = "\n".join(lines[i:i+chunk_size])
        story.append(Preformatted(chunk, code_style))
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def _convert_docx_fallback_reportlab(docx_path_or_stream):
    """Fallback puro en Python usando python-docx + ReportLab incluyendo imágenes."""
    if docx is None:
        raise ImportError("La librería python-docx no está instalada.")
        
    doc = docx.Document(docx_path_or_stream)
    buffer = io.BytesIO()
    pdf_doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle('DocNormal', parent=styles['Normal'], fontSize=10, leading=14, textColor=colors.HexColor('#334155'))
    h1_style = ParagraphStyle('DocH1', parent=styles['Heading1'], fontSize=16, leading=20, textColor=colors.HexColor('#1E293B'), spaceAfter=8)
    h2_style = ParagraphStyle('DocH2', parent=styles['Heading2'], fontSize=13, leading=17, textColor=colors.HexColor('#475569'), spaceAfter=6)
    
    story = []
    
    for p in doc.paragraphs:
        # Extraer imágenes incrustadas en el párrafo
        embed_ids = p._element.xpath('.//a:blip/@r:embed')
        if embed_ids:
            for embed_id in embed_ids:
                try:
                    if embed_id in doc.part.related_parts:
                        img_part = doc.part.related_parts[embed_id]
                        img_bytes = img_part.blob
                        img_io = io.BytesIO(img_bytes)
                        pil_img = Image.open(img_io)
                        w, h = pil_img.size
                        max_w = 480
                        if w > max_w:
                            h = int(h * (max_w / w))
                            w = max_w
                        img_io.seek(0)
                        story.append(RLImage(img_io, width=w, height=h))
                        story.append(Spacer(1, 6))
                except Exception:
                    pass

        text = p.text.strip()
        if text:
            safe_text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            if p.style and p.style.name and p.style.name.startswith('Heading 1'):
                story.append(Paragraph(safe_text, h1_style))
            elif p.style and p.style.name and p.style.name.startswith('Heading 2'):
                story.append(Paragraph(safe_text, h2_style))
            else:
                story.append(Paragraph(safe_text, normal_style))
            story.append(Spacer(1, 4))
        elif not embed_ids:
            story.append(Spacer(1, 6))
        
    for t in doc.tables:
        table_data = []
        for row in t.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                row_data.append(Paragraph(cell_text, normal_style))
            table_data.append(row_data)
        if table_data:
            t_style = TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F1F5F9')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.HexColor('#0F172A')),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
                ('TOPPADDING', (0,0), (-1,-1), 4),
                ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ])
            story.append(Table(table_data, style=t_style))
            story.append(Spacer(1, 10))
            
    if not story:
        story.append(Paragraph("(Documento Word sin contenido de texto ni imágenes)", normal_style))

    pdf_doc.build(story)
    buffer.seek(0)
    return buffer

def convert_docx_to_pdf(file_stream):
    """
    Convierte un documento Word (.docx) a PDF.
    Usa LibreOffice o Microsoft Word (win32com) si están disponibles para obtener fidelidad exacta (100%).
    De lo contrario, recurre al motor interno basado en ReportLab y python-docx con extracción de imágenes.
    """
    # Guardar stream en archivo temporal
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp_in:
        if hasattr(file_stream, 'read'):
            content = file_stream.read()
            if hasattr(file_stream, 'seek'):
                file_stream.seek(0)
            tmp_in.write(content)
        else:
            tmp_in.write(file_stream)
        tmp_in_path = tmp_in.name

    tmp_dir = tempfile.mkdtemp()
    output_pdf_path = os.path.join(tmp_dir, "output.pdf")

    try:
        # Intento 1: LibreOffice (Cross-platform / Hugging Face Spaces)
        if _convert_with_libreoffice(tmp_in_path, tmp_dir):
            if os.path.exists(output_pdf_path):
                with open(output_pdf_path, 'rb') as f:
                    return io.BytesIO(f.read())

        # Intento 2: MS Word via COM Automation (Windows local)
        if os.name == 'nt' and _convert_docx_with_win32com(tmp_in_path, output_pdf_path):
            if os.path.exists(output_pdf_path):
                with open(output_pdf_path, 'rb') as f:
                    return io.BytesIO(f.read())

        # Intento 3: Fallback en Python puro (python-docx + ReportLab con imágenes)
        return _convert_docx_fallback_reportlab(tmp_in_path)

    finally:
        if os.path.exists(tmp_in_path):
            try:
                os.remove(tmp_in_path)
            except Exception:
                pass
        if os.path.exists(tmp_dir):
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception:
                pass

def convert_xlsx_to_pdf(file_stream):
    """Convierte un libro de Excel (.xlsx) a PDF con alta fidelidad si LibreOffice está presente."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp_in:
        if hasattr(file_stream, 'read'):
            content = file_stream.read()
            if hasattr(file_stream, 'seek'):
                file_stream.seek(0)
            tmp_in.write(content)
        else:
            tmp_in.write(file_stream)
        tmp_in_path = tmp_in.name

    tmp_dir = tempfile.mkdtemp()
    output_pdf_path = os.path.join(tmp_dir, "output.pdf")

    try:
        if _convert_with_libreoffice(tmp_in_path, tmp_dir):
            if os.path.exists(output_pdf_path):
                with open(output_pdf_path, 'rb') as f:
                    return io.BytesIO(f.read())
    finally:
        if os.path.exists(tmp_in_path):
            try:
                os.remove(tmp_in_path)
            except Exception:
                pass
        if os.path.exists(tmp_dir):
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception:
                pass

    if openpyxl is None:
        raise ImportError("La librería openpyxl no está instalada.")
        
    wb = openpyxl.load_workbook(file_stream, data_only=True)
    buffer = io.BytesIO()
    pdf_doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(A4),
        rightMargin=20,
        leftMargin=20,
        topMargin=20,
        bottomMargin=20
    )
    styles = getSampleStyleSheet()
    
    sheet_title_style = ParagraphStyle('SheetTitle', parent=styles['Heading2'], fontSize=13, leading=16, textColor=colors.HexColor('#1E3A8A'), spaceAfter=8)
    cell_style = ParagraphStyle('Cell', parent=styles['Normal'], fontSize=8, leading=10)
    header_style = ParagraphStyle('HeaderCell', parent=styles['Normal'], fontSize=8, leading=10, fontName='Helvetica-Bold', textColor=colors.HexColor('#0F172A'))
    
    story = []
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        story.append(Paragraph(f"📊 Hoja: {sheet_name}", sheet_title_style))
        
        table_data = []
        for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
            if not any(row):
                continue
            row_data = []
            for cell in row:
                val = "" if cell is None else str(cell).strip()
                val_safe = val.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                if r_idx == 0:
                    row_data.append(Paragraph(val_safe, header_style))
                else:
                    row_data.append(Paragraph(val_safe, cell_style))
            table_data.append(row_data)
            
        if table_data:
            t_style = TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#E2E8F0')),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#94A3B8')),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('TOPPADDING', (0,0), (-1,-1), 3),
                ('BOTTOMPADDING', (0,0), (-1,-1), 3),
            ])
            story.append(Table(table_data, style=t_style))
            story.append(Spacer(1, 15))
            
    if not story:
        story.append(Paragraph("(Libro de Excel sin contenido)", cell_style))

    pdf_doc.build(story)
    buffer.seek(0)
    return buffer

def convert_pptx_to_pdf(file_stream):
    """Convierte una presentación PowerPoint (.pptx) a PDF con alta fidelidad si LibreOffice está presente."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_in:
        if hasattr(file_stream, 'read'):
            content = file_stream.read()
            if hasattr(file_stream, 'seek'):
                file_stream.seek(0)
            tmp_in.write(content)
        else:
            tmp_in.write(file_stream)
        tmp_in_path = tmp_in.name

    tmp_dir = tempfile.mkdtemp()
    output_pdf_path = os.path.join(tmp_dir, "output.pdf")

    try:
        if _convert_with_libreoffice(tmp_in_path, tmp_dir):
            if os.path.exists(output_pdf_path):
                with open(output_pdf_path, 'rb') as f:
                    return io.BytesIO(f.read())
    finally:
        if os.path.exists(tmp_in_path):
            try:
                os.remove(tmp_in_path)
            except Exception:
                pass
        if os.path.exists(tmp_dir):
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception:
                pass

    if pptx is None:
        raise ImportError("La librería python-pptx no está instalada.")
        
    prs = pptx.Presentation(file_stream)
    buffer = io.BytesIO()
    pdf_doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(A4),
        rightMargin=30,
        leftMargin=30,
        topMargin=30,
        bottomMargin=30
    )
    styles = getSampleStyleSheet()
    
    slide_title_style = ParagraphStyle('SlideTitle', parent=styles['Heading1'], fontSize=16, leading=20, textColor=colors.HexColor('#2563EB'), spaceAfter=10)
    body_style = ParagraphStyle('SlideBody', parent=styles['Normal'], fontSize=11, leading=15, textColor=colors.HexColor('#1F2937'))
    
    story = []
    
    for idx, slide in enumerate(prs.slides):
        story.append(Paragraph(f"🖥️ Diapositiva {idx + 1}", slide_title_style))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        safe_text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(safe_text, body_style))
                        story.append(Spacer(1, 4))
        story.append(PageBreak())
        
    if not story:
        story.append(Paragraph("(Presentación sin diapositivas)", body_style))

    pdf_doc.build(story)
    buffer.seek(0)
    return buffer

def convert_file_to_pdf(file_obj, filename=None):
    """
    Función principal que detecta el tipo de archivo y retorna un io.BytesIO con el PDF generado.
    Si el archivo ya es PDF, retorna sus bytes.
    """
    if filename is None:
        filename = getattr(file_obj, 'name', 'documento')
        
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == '.pdf':
        if hasattr(file_obj, 'read'):
            content = file_obj.read()
            if hasattr(file_obj, 'seek'):
                file_obj.seek(0)
            return io.BytesIO(content)
        return io.BytesIO(file_obj)
        
    elif ext in IMAGE_EXTENSIONS:
        return convert_image_to_pdf(file_obj)
        
    elif ext in WORD_EXTENSIONS:
        return convert_docx_to_pdf(file_obj)
        
    elif ext in EXCEL_EXTENSIONS:
        return convert_xlsx_to_pdf(file_obj)
        
    elif ext in PPT_EXTENSIONS:
        return convert_pptx_to_pdf(file_obj)
        
    elif ext in TEXT_EXTENSIONS:
        if hasattr(file_obj, 'read'):
            content = file_obj.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='replace')
            if hasattr(file_obj, 'seek'):
                file_obj.seek(0)
        else:
            content = str(file_obj)
        return convert_text_to_pdf(content, filename=filename)
        
    else:
        # Intento genérico de leer como texto plano
        try:
            if hasattr(file_obj, 'read'):
                content = file_obj.read()
                if isinstance(content, bytes):
                    content = content.decode('utf-8', errors='replace')
                if hasattr(file_obj, 'seek'):
                    file_obj.seek(0)
            else:
                content = str(file_obj)
            return convert_text_to_pdf(content, filename=filename)
        except Exception as e:
            raise ValueError(f"Formato de archivo '{ext}' no soportado para conversión: {e}")

